"""PPT 文档生成 Helper 函数
封装 python-pptx 常用操作，内置专业配色方案
"""

import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# 专业配色方案
if HAS_PPTX:
    PPT_COLOR_SCHEMES = {
        "ocean": {
            "primary": RGBColor(0x06, 0x5A, 0x82),
            "secondary": RGBColor(0x1C, 0x72, 0x93),
            "accent": RGBColor(0x21, 0x29, 0x5C),
        },
        "midnight": {
            "primary": RGBColor(0x1E, 0x27, 0x61),
            "secondary": RGBColor(0xCA, 0xDC, 0xFC),
            "accent": RGBColor(0xFF, 0xFF, 0xFF),
        },
        "forest": {
            "primary": RGBColor(0x2C, 0x5F, 0x2D),
            "secondary": RGBColor(0x97, 0xBC, 0x62),
            "accent": RGBColor(0xF5, 0xF5, 0xF5),
        },
        "coral": {
            "primary": RGBColor(0xF9, 0x61, 0x67),
            "secondary": RGBColor(0xF9, 0xE7, 0x95),
            "accent": RGBColor(0x2F, 0x3C, 0x7E),
        },
        "charcoal": {
            "primary": RGBColor(0x36, 0x45, 0x4F),
            "secondary": RGBColor(0xF2, 0xF2, 0xF2),
            "accent": RGBColor(0x21, 0x21, 0x21),
        },
    }
else:
    PPT_COLOR_SCHEMES = {}


def create_ppt_doc(title=None, author="", color_scheme="ocean"):
    """创建一个预配置好专业样式的 PPT 演示文稿对象

    Args:
        title: 演示文稿标题（可选）
        author: 文档作者
        color_scheme: 配色方案名称 (ocean/midnight/forest/coral/charcoal)

    Returns:
        Presentation: 预配置好的 python-pptx Presentation 对象

    示例:
        prs = create_ppt_doc(title="项目汇报", color_scheme="ocean")
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        save_ppt_doc(prs, "项目汇报.pptx")
    """
    prs = Presentation()

    if title:
        prs.core_properties.title = title
    if author:
        prs.core_properties.author = author

    return prs


def save_ppt_doc(prs, filename, working_dir=""):
    """保存 PPT 演示文稿到工作目录

    Args:
        prs: python-pptx Presentation 对象
        filename: 文件名（如 "演示.pptx"）
        working_dir: 工作目录路径

    Returns:
        str: 保存的文件绝对路径
    """
    output_path = os.path.join(working_dir, filename) if working_dir else filename
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)

    return output_path
