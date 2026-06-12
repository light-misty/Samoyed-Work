"""PPT Handler - PowerPoint document processor
Based on python-pptx, supports read/convert/analyze operations.
精简版：仅支持 read/convert/analyze 操作
"""

import os
import logging

from pptx import Presentation


class PptHandler:
    """PowerPoint (.pptx) 文档处理器（精简版，仅支持 read/convert/analyze）"""

    logger = logging.getLogger(__name__)

    def read(self, params: dict) -> dict:
        path = params.get("path", "")
        if not path:
            self.logger.error("read: missing path")
            return {"error": "missing path"}
        if not os.path.exists(path):
            raise FileNotFoundError(path)
        self.logger.info("read: loading PPT %s", path)
        prs = Presentation(path)
        slides = []
        for slide in prs.slides:
            info = {"shapes": []}
            for shape in slide.shapes:
                si = {"name": shape.name, "type": str(shape.shape_type)}
                if shape.has_text_frame:
                    si["text"] = shape.text
                info["shapes"].append(si)
            slides.append(info)
        self.logger.info("read: done, %d slides", len(slides))
        return {"slides": slides, "slide_count": len(slides)}

    def convert(self, params: dict) -> dict:
        path = params.get("path", "")
        output_path = params.get("output_path", "")
        target = params.get("format", "pdf")
        if not path:
            self.logger.error("convert: missing path")
            return {"error": "missing path"}
        if not os.path.exists(path):
            raise FileNotFoundError(path)
        self.logger.info("convert: %s -> %s", path, target)
        if target == "pdf":
            out = output_path or os.path.splitext(path)[0] + ".pdf"
            return {"path": out, "format": target, "message": "PPT to PDF requires LibreOffice"}
        return {"error": "unsupported format: " + target}

    def analyze(self, params: dict) -> dict:
        path = params.get("path", "")
        if not path:
            self.logger.error("analyze: missing path")
            return {"error": "missing path"}
        if not os.path.exists(path):
            raise FileNotFoundError(path)
        self.logger.info("analyze: loading PPT %s", path)
        prs = Presentation(path)
        ts = 0
        tts = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                ts += 1
                if shape.has_text_frame:
                    tts += 1
        self.logger.info("analyze: done, %d slides", len(prs.slides))
        return {"file_size": os.path.getsize(path), "slide_count": len(prs.slides), "total_shapes": ts, "total_text_shapes": tts}
